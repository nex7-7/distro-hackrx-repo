"""Ingestion Engine

Implements a unified file ingestion pipeline with a parser registry and a main
`ingest_from_url` function. Keeps outward behaviour (returning List[str] of
raw text chunks) consistent with the existing retrieval pipeline.

Features:
 - Download remote file (or accept local path) to temp storage
 - Detect extension and route to appropriate parser via registry
 - Handle ZIP archives recursively
 - Parse PDFs with fast PyMuPDF pathway
 - Use unstructured's partitioning directly (if installed) for generic documents
     / spreadsheets / images (best‑effort OCR)
 - Graceful fallbacks & detailed logging
"""
from __future__ import annotations

import os
import uuid
import time
import zipfile
import tempfile
import requests
from typing import Callable, Dict, List, Optional
from urllib.parse import urlparse, unquote

from fastapi import HTTPException

from components.utils.logger import log_service_event, log_error

# Optional direct import of unstructured partition function; engine degrades gracefully if absent.
try:  # pragma: no cover
    from unstructured.partition.auto import partition as _unstructured_partition  # type: ignore
    _UNSTRUCTURED_AVAILABLE = True
except Exception:  # pragma: no cover
    _unstructured_partition = None  # type: ignore
    _UNSTRUCTURED_AVAILABLE = False

log_service_event(
    "unstructured_availability",
    f"unstructured partition module {'available' if _UNSTRUCTURED_AVAILABLE else 'unavailable'}",
    {"available": _UNSTRUCTURED_AVAILABLE}
)

# ------------------------------ Parser Functions ------------------------------


def parse_standard_document(file_path: str, source_url: str) -> List[str]:
    """Parse a standard text / office document into list[str].

    Uses unstructured.process_document to obtain chunk groups; flattens each
    chunk (list of element dicts) into a single string joined by double newlines.
    Falls back to naive read for .txt if unstructured unavailable.
    """
    start = time.time()
    ext = os.path.splitext(file_path.lower())[1]
    chunks: List[str] = []
    if not _UNSTRUCTURED_AVAILABLE:
        if ext == ".txt":  # basic text fallback
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
            chunks = [seg.strip() for seg in text.split('\n\n') if seg.strip()]
            log_service_event("plain_text_fallback", "Processed TXT without unstructured", {
                "file_path": file_path,
                "chunks": len(chunks)
            })
        else:
            log_service_event("unstructured_unavailable_skip", "Skipping non-txt file (unstructured missing)", {
                "file_path": file_path,
                "ext": ext
            })
            return []
    else:
        try:
            # Lightweight internal grouping to mimic previous behaviour
            elements = _unstructured_partition(
                filename=file_path, languages=["eng"]) or []
            if not elements:
                try:
                    elements = _unstructured_partition(
                        filename=file_path, languages=None) or []
                except Exception:
                    pass
            groups = []
            current = []
            current_words = 0
            for el in elements:
                text = str(el).strip()
                if not text:
                    continue
                words = text.split()
                if current_words + len(words) > 500 and current:
                    groups.append(current)
                    current = []
                    current_words = 0
                # Capture dict form if available
                try:
                    d = el.to_dict()  # type: ignore[attr-defined]
                    if 'text' not in d:
                        d['text'] = text
                except Exception:
                    d = {'text': text}
                current.append(d)
                current_words += len(words)
            if current:
                groups.append(current)
            for group in groups:
                texts = []
                for el in group:
                    txt = el.get('text') if isinstance(el, dict) else None
                    if not txt:
                        # element objects handled by str(el)
                        try:
                            txt = el.get('text')  # type: ignore
                        except Exception:
                            txt = str(el)
                    if txt and txt.strip():
                        texts.append(txt.strip())
                if texts:
                    chunks.append("\n\n".join(texts))
        except Exception as e:  # pragma: no cover
            log_error("standard_doc_parse_failed", {
                "file_path": file_path,
                "error": str(e)
            })
            raise HTTPException(
                status_code=422, detail=f"Failed to parse document: {e}")

    log_service_event("standard_doc_parsed", "Parsed standard document", {
        "file_path": file_path,
        "source_url": source_url,
        "ext": ext,
        "chunks": len(chunks),
        "time_seconds": time.time() - start
    })
    return chunks


def parse_spreadsheet(file_path: str, source_url: str) -> List[str]:
    """Parse spreadsheet (.xlsx) into list[str] using pandas.

    Each sheet will be processed separately with preserved structure for better readability.
    Tables are formatted to maintain alignment and preserve relationships between data cells.
    """
    try:
        import pandas as pd
    except ImportError:
        log_error("pandas_import_failed", {
                  "file_path": file_path, "error": "pandas not installed"})
        # Fall back to standard parser
        return parse_standard_document(file_path, source_url)

    start_time = time.time()
    chunks = []

    try:
        # Read all sheets from the Excel file
        xl = pd.ExcelFile(file_path)
        sheet_names = xl.sheet_names

        for sheet_name in sheet_names:
            # Read each sheet into a dataframe
            df = pd.read_excel(file_path, sheet_name=sheet_name)

            # Skip empty sheets
            if df.empty:
                continue

            # Create a metadata header for the sheet
            sheet_header = f"TABLE: {sheet_name}"
            sheet_info = f"Rows: {df.shape[0]}, Columns: {df.shape[1]}"
            metadata = f"{sheet_header}\n{sheet_info}\n{'=' * 50}"

            # Preserve structure by converting to markdown table format
            # First, convert any numeric indices to strings
            df.columns = df.columns.astype(str)

            # Format column headers with proper alignment
            headers = " | ".join([str(col).strip() for col in df.columns])
            separator = "-|-".join(["-" * len(str(col)) for col in df.columns])

            # Build the table rows with proper alignment
            rows = []
            for _, row in df.iterrows():
                formatted_row = " | ".join(
                    [str(val).strip() if pd.notna(val) else "" for val in row])
                rows.append(formatted_row)

            # Combine all parts into a well-structured table
            table_content = f"{headers}\n{separator}\n" + "\n".join(rows)
            structured_table = f"{metadata}\n\n{table_content}"

            # Add the formatted table as a chunk
            chunks.append(structured_table)

            # For large tables, also create chunks by row groups to improve search relevance
            if df.shape[0] > 10:  # More than 10 rows
                # Adaptive chunk size
                row_chunk_size = min(10, max(5, df.shape[0] // 5))

                for i in range(0, df.shape[0], row_chunk_size):
                    end_idx = min(i + row_chunk_size, df.shape[0])
                    row_group = df.iloc[i:end_idx]

                    # Create a sub-table with the same format
                    sub_headers = " | ".join(
                        [str(col).strip() for col in row_group.columns])
                    sub_separator = "-|-".join(["-" * len(str(col))
                                               for col in row_group.columns])

                    sub_rows = []
                    for _, row in row_group.iterrows():
                        formatted_row = " | ".join(
                            [str(val).strip() if pd.notna(val) else "" for val in row])
                        sub_rows.append(formatted_row)

                    sub_table = f"{sub_headers}\n{sub_separator}\n" + \
                        "\n".join(sub_rows)
                    sub_chunk = f"{sheet_header} - Rows {i+1}-{end_idx}\n{'=' * 30}\n\n{sub_table}"
                    chunks.append(sub_chunk)

    except Exception as e:
        log_error("excel_pandas_parse_failed", {
                  "file_path": file_path, "error": str(e)})
        # Fall back to standard parser
        return parse_standard_document(file_path, source_url)

    # Filter out very short chunks
    chunks = [chunk for chunk in chunks if len(chunk) > 50]

    parsing_time = time.time() - start_time
    log_service_event("excel_pandas_parsed", "Parsed Excel with pandas", {
        "file_path": file_path,
        "source_url": source_url,
        "sheets_found": len(sheet_names),
        "chunks": len(chunks),
        "time_seconds": parsing_time
    })

    return chunks


def parse_image_with_ocr(file_path: str, source_url: str) -> List[str]:
    """Parse image using unstructured (which may invoke OCR if deps present).

    If OCR backend not available, returns empty list (logged) instead of failing.
    """
    try:
        return parse_standard_document(file_path, source_url)
    except HTTPException as e:
        log_service_event("image_ocr_unavailable", "Image OCR unavailable or failed", {
            "file_path": file_path,
            "source_url": source_url,
            "detail": str(e)
        })
        return []


def parse_pdf_file(file_path: str, source_url: str) -> List[str]:
    """Parse PDF using PyMuPDF directly (inlined) and return chunk list.

    Mirrors previous parse_pdf_with_pymupdf behaviour to keep identical
    splitting heuristic (paragraphs / long paragraph sentence splitting).
    """
    import fitz  # PyMuPDF
    start_time = time.time()
    document_id = str(uuid.uuid4())
    filename = os.path.basename(file_path)

    log_service_event("primary_parsing_start", "Starting primary PDF parsing with PyMuPDF", {
        "document_id": document_id,
        "filename": filename,
        "parsing_method": "PyMuPDF"
    })

    try:
        doc = fitz.open(file_path)
        page_texts: List[str] = []
        for pno in range(len(doc)):
            page = doc.load_page(pno)
            txt = page.get_text("text")
            if txt.strip():
                page_texts.append(txt)
        doc.close()

        if not page_texts:
            log_error("pdf_extraction_empty", {
                      "document_id": document_id, "filename": filename})
            return []

        chunks: List[str] = []
        for page_text in page_texts:
            paragraphs = page_text.split("\n\n")
            for para in paragraphs:
                if not para.strip():
                    continue
                if len(para) > 1000:
                    sentences = para.split(". ")
                    current = ""
                    for s in sentences:
                        if len(current) + len(s) < 1000:
                            current += s + ". "
                        else:
                            if current:
                                chunks.append(current.strip())
                            current = s + ". "
                    if current:
                        chunks.append(current.strip())
                else:
                    chunks.append(para.strip())

        # filter short
        chunks = [c for c in chunks if len(c.strip()) > 50]
        parsing_time = time.time() - start_time
        log_service_event("primary_parsing_complete", "Successfully parsed PDF with PyMuPDF", {
            "document_id": document_id,
            "filename": filename,
            "chunks_extracted": len(chunks),
            "total_chars": sum(len(c) for c in chunks),
            "avg_chunk_length": (sum(len(c) for c in chunks) / len(chunks)) if chunks else 0,
            "parsing_time_seconds": parsing_time
        })
        return chunks
    except Exception as e:  # pragma: no cover
        log_error("primary_parsing_failed", {
                  "document_id": document_id, "filename": filename, "error": str(e)})
        return []


# --- PPTX Parser using Docling and EasyOCR for images ---
def parse_pptx_file(file_path: str, source_url: str) -> List[str]:
    """Parse PPTX file and return list of text chunks using Docling.

    If images are detected (shown as <!-- image --> tags), extract them
    and use EasyOCR to perform OCR on the images for better text extraction.
    """
    try:
        from docling.document_converter import DocumentConverter
    except ImportError:
        log_error("docling_import_failed", {
            "file_path": file_path,
            "error": "Required module not installed (docling). Try: pip install docling docling-core"
        })
        # Fall back to the old parser if available
        try:
            from pptx import Presentation
            log_service_event("docling_fallback", "Falling back to python-pptx parser", {
                "file_path": file_path
            })
            return _legacy_parse_pptx_file(file_path, source_url)
        except ImportError:
            return []

    start_time = time.time()
    chunks = []

    try:
        # Use Docling to convert the PPTX file
        converter = DocumentConverter()
        result = converter.convert(file_path)
        document = result.document

        # Get markdown representation
        markdown_content = document.export_to_markdown()

        # Check if markdown contains image placeholders
        contains_images = "<!-- image -->" in markdown_content

        # Split markdown by headings to get slide chunks
        slide_chunks = []
        current_slide = []
        slide_count = 0
        lines = markdown_content.split('\n')
        for line in lines:
            if line.startswith('# ') or line.startswith('## '):
                if current_slide:
                    slide_count += 1
                    slide_text = '\n\n'.join(current_slide)
                    if not slide_text.startswith('Slide'):
                        slide_text = f"Slide {slide_count}: {current_slide[0]}\n\n" + slide_text
                    slide_chunks.append(slide_text)
                    current_slide = []
            if line.strip():
                current_slide.append(line)
        if current_slide:
            slide_count += 1
            slide_text = '\n\n'.join(current_slide)
            if not slide_text.startswith('Slide'):
                slide_text = f"Slide {slide_count}: {current_slide[0]}\n\n" + slide_text
            slide_chunks.append(slide_text)

        chunks = slide_chunks

        # If images were detected, extract them from PPTX and perform OCR
        if contains_images:
            try:
                log_service_event("image_extraction_start", "Images detected in PPTX, extracting with python-pptx", {
                    "file_path": file_path,
                })

                # Extract images using python-pptx
                from pptx import Presentation
                import tempfile
                import os

                # Create temp directory for extracted images
                images_temp_dir = tempfile.mkdtemp(prefix="pptx_images_")
                prs = Presentation(file_path)

                image_ocr_results = []
                img_count = 0

                # Extract images from each slide
                for i, slide in enumerate(prs.slides):
                    slide_num = i + 1
                    slide_images = []

                    for shape in slide.shapes:
                        if shape.shape_type == 13:  # 13 is the enum value for pictures
                            img_count += 1
                            try:
                                # Extract image to temp file
                                image = shape.image
                                image_bytes = image.blob
                                img_ext = image.ext.lower() if hasattr(image, 'ext') else '.png'

                                # Ensure we handle common image extensions correctly
                                if img_ext not in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff']:
                                    img_ext = '.png'  # Default to PNG for unknown formats

                                img_path = os.path.join(
                                    images_temp_dir, f"slide_{slide_num}_img_{img_count}{img_ext}")

                                with open(img_path, 'wb') as img_file:
                                    img_file.write(image_bytes)

                                # Use EasyOCR on the image
                                try:
                                    import easyocr
                                    import numpy as np
                                    import cv2

                                    # Initialize EasyOCR reader once (lazy loading)
                                    if 'ocr_reader' not in locals():
                                        log_service_event(
                                            "easyocr_init", "Initializing EasyOCR", {})
                                        ocr_reader = easyocr.Reader(
                                            ['en'], gpu=False)

                                    # Read image and perform OCR
                                    try:
                                        # For JPEG/JPG files, we'll try using PIL first if cv2 fails
                                        img = cv2.imread(img_path)

                                        if img is None and (img_path.lower().endswith('.jpg') or img_path.lower().endswith('.jpeg')):
                                            from PIL import Image
                                            import numpy as np
                                            pil_img = Image.open(img_path)
                                            img = np.array(
                                                pil_img.convert('RGB'))
                                            # Convert RGB to BGR for OpenCV
                                            img = img[:, :, ::-1].copy()

                                        if img is not None:
                                            ocr_results = ocr_reader.readtext(
                                                img)

                                            if ocr_results:
                                                # Extract text from OCR results
                                                texts = [
                                                    text for _, text, conf in ocr_results if conf > 0.5]
                                                if texts:
                                                    ocr_text = " ".join(texts)
                                                    image_ocr_results.append(
                                                        f"Slide {slide_num} Image {img_count} OCR: {ocr_text}")
                                                    log_service_event("easyocr_success", "Successfully extracted text with EasyOCR", {
                                                        "slide": slide_num,
                                                        "image": img_count,
                                                        "format": img_ext,
                                                        "text_length": len(ocr_text)
                                                    })
                                    except Exception as img_read_err:
                                        log_error("image_read_failed", {
                                            "file_path": img_path,
                                            "format": img_ext,
                                            "error": str(img_read_err)
                                        })
                                except Exception as ocr_err:
                                    log_error("easyocr_failed", {
                                        "file_path": img_path,
                                        "slide": slide_num,
                                        "error": str(ocr_err)
                                    })
                            except Exception as img_err:
                                log_error("image_extraction_failed", {
                                    "slide": slide_num,
                                    "error": str(img_err)
                                })

                # Add OCR results to chunks
                if image_ocr_results:
                    for i, ocr_result in enumerate(image_ocr_results):
                        # Find which slide this image belongs to based on slide number in the OCR result
                        slide_num = int(ocr_result.split(
                            "Slide ")[1].split(" ")[0])

                        # Try to find the corresponding slide chunk
                        for j, chunk in enumerate(chunks):
                            if chunk.startswith(f"Slide {slide_num}:") or chunk.startswith(f"Slide {slide_num}\n"):
                                # Append OCR result to the slide content
                                chunks[j] = chunks[j] + f"\n\n{ocr_result}"
                                break

                # Clean up temp directory
                try:
                    import shutil
                    shutil.rmtree(images_temp_dir)
                except Exception:
                    pass

            except Exception as extract_err:
                log_error("pptx_image_extraction_failed", {
                    "file_path": file_path,
                    "error": str(extract_err)
                })

    except Exception as e:
        log_error("docling_pptx_parse_failed", {
            "file_path": file_path,
            "error": str(e)
        })
        # Fall back to legacy parser
        try:
            log_service_event("docling_error_fallback", "Falling back to python-pptx parser after Docling error", {
                "file_path": file_path,
                "error": str(e)
            })
            return _legacy_parse_pptx_file(file_path, source_url)
        except Exception:
            return []

    parsing_time = time.time() - start_time
    log_service_event("pptx_parsed_with_docling", "Parsed PPTX file with Docling", {
        "file_path": file_path,
        "source_url": source_url,
        "slides_processed": len(chunks),
        "contains_images": contains_images,
        "parsing_time_seconds": parsing_time
    })
    return chunks


# Legacy PPTX parser as fallback
def _legacy_parse_pptx_file(file_path: str, source_url: str) -> List[str]:
    """Legacy PPTX parser using python-pptx as a fallback if Docling is unavailable."""
    try:
        from pptx import Presentation
        import io
        import os
        from PIL import Image
    except ImportError:
        log_error("pptx_import_failed", {
                  "file_path": file_path, "error": "Required modules not installed (python-pptx, pillow)"})
        return []

    chunks = []
    images_temp_dir = None

    try:
        # Create temp directory for extracted images
        images_temp_dir = tempfile.mkdtemp(prefix="pptx_images_")
        prs = Presentation(file_path)

        for i, slide in enumerate(prs.slides):
            slide_num = i + 1
            slide_content = [f"Slide {slide_num}"]

            # Get slide title if available
            title = None
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip() and hasattr(shape, 'is_title') and shape.is_title:
                    title = shape.text.strip()
                    break

            if title:
                slide_content[0] = f"Slide {slide_num}: {title}"

            # Extract text from all text shapes
            text_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text.strip())

            # Combine all text content
            if text_content:
                slide_content.append("Text Content:")
                slide_content.append("\n".join(text_content))

            # Extract and process images if available
            image_content = []
            img_count = 0

            for shape in slide.shapes:
                if shape.shape_type == 13:  # 13 is the enum value for pictures
                    img_count += 1
                    try:
                        # Extract image to temp file
                        image = shape.image
                        image_bytes = image.blob
                        img_ext = image.ext.lower() if hasattr(image, 'ext') else '.png'
                        img_path = os.path.join(
                            images_temp_dir, f"slide_{slide_num}_img_{img_count}{img_ext}")

                        with open(img_path, 'wb') as img_file:
                            img_file.write(image_bytes)

                        # Use OCR on the image if unstructured is available
                        if _UNSTRUCTURED_AVAILABLE:
                            img_text = []
                            try:
                                elements = _unstructured_partition(
                                    filename=img_path) or []
                                for el in elements:
                                    text = str(el).strip()
                                    if text:
                                        img_text.append(text)
                            except Exception as ocr_err:
                                log_error("image_ocr_failed", {
                                    "file_path": img_path,
                                    "slide": slide_num,
                                    "error": str(ocr_err)
                                })

                            if img_text:
                                image_content.append(
                                    f"Image {img_count} content: {' '.join(img_text)}")
                    except Exception as img_err:
                        log_error("image_extraction_failed", {
                            "slide": slide_num,
                            "error": str(img_err)
                        })

            # Add image content if available
            if image_content:
                slide_content.append("\nImage Content:")
                slide_content.extend(image_content)

            # Add this slide's content as a chunk
            if len(slide_content) > 1:  # More than just the slide header
                chunks.append("\n\n".join(slide_content))

    except Exception as e:
        log_error("pptx_parse_failed", {
                  "file_path": file_path, "error": str(e)})
        return []
    finally:
        # Clean up temp directory
        if images_temp_dir and os.path.exists(images_temp_dir):
            try:
                import shutil
                shutil.rmtree(images_temp_dir)
            except Exception:
                pass

    log_service_event("pptx_parsed_legacy", "Parsed PPTX file with legacy parser", {
        "file_path": file_path,
        "source_url": source_url,
        "slides_processed": len(prs.slides) if 'prs' in locals() else 0,
        "chunks": len(chunks)
    })
    return chunks


# ------------------------------ File Security Configuration ------------------------------

# Maximum file size allowed (50MB in bytes)
MAX_FILE_SIZE_BYTES = 1000 * 1024 * 1024  # 1GB

# Blacklisted file extensions that should not be processed
BLACKLISTED_EXTENSIONS = {
    "bin", "exe", "dll", "so", "dylib", "app", "deb", "rpm",
    "msi", "dmg", "pkg", "run", "tar", "gz", "bz2", "xz", "7z", "rar"
}

# Maximum ZIP extraction depth to prevent recursive ZIP bombs
MAX_ZIP_DEPTH = 1

# ------------------------------ Security Functions ------------------------------

def _check_file_security(file_path: str, current_zip_depth: int = 0) -> bool:
    """
    Check if a file is safe to process based on security constraints.
    
    Parameters:
        file_path (str): Path to the file to check
        current_zip_depth (int): Current depth of ZIP extraction (0 = not in ZIP)
    
    Returns:
        bool: True if file is safe to process, False otherwise
    
    Raises:
        HTTPException: If file violates security constraints
    """
    # Check file size
    try:
        file_size = os.path.getsize(file_path)
        if file_size > MAX_FILE_SIZE_BYTES:
            log_error("file_too_large", {
                "file_path": file_path,
                "size_bytes": file_size,
                "max_size_bytes": MAX_FILE_SIZE_BYTES,
                "size_mb": file_size / (1024 * 1024)
            })
            raise HTTPException(
                status_code=413,
                detail=f"File too large: {file_size / (1024 * 1024):.1f}MB. Maximum allowed: {MAX_FILE_SIZE_BYTES / (1024 * 1024)}MB"
            )
    except OSError as e:
        log_error("file_size_check_failed", {
            "file_path": file_path,
            "error": str(e)
        })
        raise HTTPException(
            status_code=400,
            detail=f"Cannot access file: {e}"
        )
    
    # Check file extension against blacklist
    ext = os.path.splitext(file_path)[1].lower().lstrip('.')
    if ext in BLACKLISTED_EXTENSIONS:
        log_error("blacklisted_file_type", {
            "file_path": file_path,
            "extension": ext,
            "blacklisted_extensions": list(BLACKLISTED_EXTENSIONS)
        })
        raise HTTPException(
            status_code=415,
            detail=f"File type '.{ext}' is not allowed for security reasons"
        )
    
    # Check ZIP depth to prevent recursive ZIP bombs
    if ext == "zip" and current_zip_depth >= MAX_ZIP_DEPTH:
        log_error("recursive_zip_blocked", {
            "file_path": file_path,
            "current_depth": current_zip_depth,
            "max_depth": MAX_ZIP_DEPTH
        })
        raise HTTPException(
            status_code=422,
            detail=f"Recursive ZIP files not allowed. Maximum extraction depth: {MAX_ZIP_DEPTH}"
        )
    
    log_service_event("file_security_check_passed", "File passed security checks", {
        "file_path": file_path,
        "extension": ext,
        "size_bytes": file_size,
        "zip_depth": current_zip_depth
    })
    
    return True

# ------------------------------ Parser Registry ------------------------------
Parser_Registry: Dict[str, Callable[[str, str], List[str]]] = {
    "xlsx": parse_spreadsheet,
    "jpg": parse_image_with_ocr,
    "jpeg": parse_image_with_ocr,
    "png": parse_image_with_ocr,
    "pdf": parse_pdf_file,
    "docx": parse_standard_document,
    "txt": parse_standard_document,
    "pptx": parse_pptx_file,
}

# ------------------------------ Helper Functions ------------------------------


def _is_url(path_or_url: str) -> bool:
    parsed = urlparse(path_or_url)
    return bool(parsed.scheme and parsed.netloc)


def _download_to_temp(url: str) -> str:
    """Download a file from URL to temporary storage with security checks.
    
    Parameters:
        url (str): URL to download from
    
    Returns:
        str: Path to the downloaded temporary file
        
    Raises:
        HTTPException: If download fails or file violates security constraints
    """
    dl_id = str(uuid.uuid4())
    log_service_event("download_start", "Starting file download", {
        "url": url,
        "download_id": dl_id
    })
    
    try:
        # Download with streaming to check size during download
        resp = requests.get(url, timeout=60, stream=True)
        resp.raise_for_status()
        
        # Check Content-Length header if available
        content_length = resp.headers.get('content-length')
        if content_length and int(content_length) > MAX_FILE_SIZE_BYTES:
            log_error("download_file_too_large", {
                "url": url,
                "content_length": int(content_length),
                "max_size_bytes": MAX_FILE_SIZE_BYTES
            })
            raise HTTPException(
                status_code=413,
                detail=f"Remote file too large: {int(content_length) / (1024 * 1024):.1f}MB. Maximum allowed: {MAX_FILE_SIZE_BYTES / (1024 * 1024)}MB"
            )
        
        # Determine file extension from URL
        name = os.path.basename(unquote(urlparse(url).path)) or f"file_{dl_id}"
        suffix = os.path.splitext(name)[1] or ''
        
        # Check if extension is blacklisted before downloading
        ext = suffix.lower().lstrip('.')
        if ext in BLACKLISTED_EXTENSIONS:
            log_error("download_blacklisted_extension", {
                "url": url,
                "extension": ext,
                "filename": name
            })
            raise HTTPException(
                status_code=415,
                detail=f"File type '.{ext}' is not allowed for security reasons"
            )
        
        # Create temporary file
        tmp_fd, tmp_path = tempfile.mkstemp(prefix="ingest_", suffix=suffix)
        
        # Download with size checking
        downloaded_bytes = 0
        chunk_size = 8192
        
        try:
            with os.fdopen(tmp_fd, 'wb') as f:
                for chunk in resp.iter_content(chunk_size=chunk_size):
                    if chunk:
                        downloaded_bytes += len(chunk)
                        
                        # Check size limit during download
                        if downloaded_bytes > MAX_FILE_SIZE_BYTES:
                            log_error("download_size_limit_exceeded", {
                                "url": url,
                                "downloaded_bytes": downloaded_bytes,
                                "max_size_bytes": MAX_FILE_SIZE_BYTES
                            })
                            raise HTTPException(
                                status_code=413,
                                detail=f"File too large during download: {downloaded_bytes / (1024 * 1024):.1f}MB. Maximum allowed: {MAX_FILE_SIZE_BYTES / (1024 * 1024)}MB"
                            )
                        
                        f.write(chunk)
        except HTTPException:
            # Clean up temp file on error
            try:
                os.unlink(tmp_path)
            except OSError:
                pass
            raise
        
    except requests.exceptions.RequestException as e:
        log_error("download_failed", {"url": url, "error": str(e)})
        raise HTTPException(
            status_code=400, detail=f"Failed to download file: {e}")
    except HTTPException:
        # Re-raise HTTP exceptions
        raise
    except Exception as e:
        log_error("download_unexpected_error", {"url": url, "error": str(e)})
        raise HTTPException(
            status_code=500, detail=f"Unexpected error during download: {e}")

    log_service_event("download_complete", "File downloaded successfully", {
        "url": url,
        "path": tmp_path,
        "bytes": downloaded_bytes,
        "download_id": dl_id
    })
    return tmp_path


def handle_zip_file(zip_path: str, source_url: str, current_zip_depth: int = 0) -> List[str]:
    """
    Handle ZIP file extraction with security checks and depth limits.
    
    Parameters:
        zip_path (str): Path to the ZIP file
        source_url (str): Original source URL for logging
        current_zip_depth (int): Current depth of ZIP extraction
    
    Returns:
        List[str]: Combined chunks from all files in the ZIP
    """
    zip_id = str(uuid.uuid4())
    start = time.time()
    
    # Check ZIP depth before processing
    if current_zip_depth >= MAX_ZIP_DEPTH:
        log_error("max_zip_depth_exceeded", {
            "zip_id": zip_id,
            "zip_path": zip_path,
            "current_depth": current_zip_depth,
            "max_depth": MAX_ZIP_DEPTH
        })
        raise HTTPException(
            status_code=422, 
            detail=f"Maximum ZIP extraction depth ({MAX_ZIP_DEPTH}) exceeded. Recursive ZIPs not allowed."
        )
    
    log_service_event("zip_extraction_start", "Extracting ZIP", {
        "zip_id": zip_id,
        "zip_path": zip_path,
        "extraction_depth": current_zip_depth
    })
    
    combined: List[str] = []
    processable_files_found = False
    
    with tempfile.TemporaryDirectory() as temp_dir:
        try:
            with zipfile.ZipFile(zip_path, 'r') as zf:
                # Check for suspicious ZIP structure before extraction
                members = zf.namelist()
                
                # Filter out suspicious files and check for depth issues
                safe_members = []
                for member in members:
                    # Skip hidden/system files
                    if any(part.startswith('.') or part.startswith('__MACOSX') for part in member.split('/')):
                        continue
                    
                    # Check if this is a nested ZIP
                    if member.lower().endswith('.zip') and current_zip_depth >= MAX_ZIP_DEPTH - 1:
                        log_service_event("nested_zip_skipped", "Skipping nested ZIP due to depth limit", {
                            "zip_id": zip_id,
                            "nested_zip": member,
                            "current_depth": current_zip_depth
                        })
                        continue
                    
                    safe_members.append(member)
                
                if not safe_members:
                    log_service_event("zip_no_processable_files", "No processable files found in ZIP", {
                        "zip_id": zip_id,
                        "total_members": len(members),
                        "safe_members": 0
                    })
                    return []
                
                # Extract only safe members
                for member in safe_members:
                    try:
                        zf.extract(member, temp_dir)
                    except Exception as e:
                        log_error("zip_member_extraction_failed", {
                            "zip_id": zip_id,
                            "member": member,
                            "error": str(e)
                        })
                        continue
                
            log_service_event("zip_extracted", "ZIP extracted successfully", {
                "zip_id": zip_id,
                "total_members": len(members),
                "safe_members": len(safe_members),
                "extraction_depth": current_zip_depth
            })
            
        except zipfile.BadZipFile as e:
            log_error("invalid_zip", {"zip_path": zip_path, "error": str(e)})
            raise HTTPException(
                status_code=422, detail=f"Invalid ZIP file: {e}")
        except Exception as e:
            log_error("zip_extraction_failed", {
                "zip_id": zip_id,
                "zip_path": zip_path,
                "error": str(e)
            })
            raise HTTPException(
                status_code=500, detail=f"Failed to extract ZIP: {e}")

        # Process extracted files
        for root, _, files in os.walk(temp_dir):
            for fname in files:
                if fname.startswith("._"):  # Skip macOS metadata files
                    continue
                
                fpath = os.path.join(root, fname)
                
                try:
                    # Check file security before processing
                    _check_file_security(fpath, current_zip_depth + 1)
                    
                    # Determine if this is a processable file type
                    ext = os.path.splitext(fname.lower())[1].lstrip('.')
                    if ext in Parser_Registry or ext == 'zip':
                        processable_files_found = True
                    
                    # Process the file
                    chunks = ingest_from_url(
                        fpath, 
                        _is_recursive=True, 
                        _source_url=source_url,
                        _current_zip_depth=current_zip_depth + 1
                    )
                    combined.extend(chunks)
                    
                except HTTPException as he:
                    # Re-raise HTTP exceptions (these are security violations)
                    if he.status_code in [413, 415, 422]:  # File too large, blacklisted, or recursive ZIP
                        raise he
                    else:
                        log_error("zip_member_security_check_failed", {
                            "zip_id": zip_id,
                            "member": fname,
                            "error": str(he)
                        })
                        continue
                except Exception as e:
                    log_error("zip_member_ingest_failed", {
                        "zip_id": zip_id,
                        "member": fname,
                        "error": str(e)
                    })
                    continue
        
        # Check if we found any processable content
        if not processable_files_found or not combined:
            log_service_event("zip_no_content_extracted", "No valid content extracted from ZIP", {
                "zip_id": zip_id,
                "processable_files_found": processable_files_found,
                "chunks_extracted": len(combined)
            })
            # Return empty list which will trigger "Files Not found" response
            return []
        
        log_service_event("zip_processing_complete", "ZIP processed successfully", {
            "zip_id": zip_id,
            "total_chunks": len(combined),
            "processable_files_found": processable_files_found,
            "elapsed": time.time() - start,
            "extraction_depth": current_zip_depth
        })
        
    return combined

# ------------------------------ Main Engine ------------------------------


def ingest_from_url(url_or_path: str, *, _is_recursive: bool = False, _source_url: Optional[str] = None, _current_zip_depth: int = 0) -> List[str]:
    """Ingest a document (remote URL or local path) and return list[str] chunks.

    Implements required steps:
      1. Download file if remote
      2. Security checks (file size, type, depth)
      3. Detect extension
      4. If ZIP -> recurse into contents (with depth limits)
      5. Route to parser via registry
      6. Execute parser & collect chunks
    
    Parameters:
        url_or_path (str): URL or local file path to process
        _is_recursive (bool): Internal flag indicating recursive call
        _source_url (Optional[str]): Original source URL for logging
        _current_zip_depth (int): Current ZIP extraction depth
    
    Returns:
        List[str]: List of text chunks extracted from the document
        
    Raises:
        HTTPException: For security violations or processing errors
    """
    original_input = url_or_path
    start = time.time()
    temp_path = None
    is_local = os.path.exists(url_or_path)
    source_url = _source_url or (url_or_path if _is_url(url_or_path) else None)

    try:
        # Step 1: Download file if remote
        if not is_local:
            temp_path = _download_to_temp(url_or_path)
        else:
            temp_path = url_or_path

        # Step 2: Security checks
        try:
            _check_file_security(temp_path, _current_zip_depth)
        except HTTPException as he:
            # For security violations, log and re-raise
            log_error("file_security_violation", {
                "input": original_input,
                "file_path": temp_path,
                "zip_depth": _current_zip_depth,
                "status_code": he.status_code,
                "detail": he.detail
            })
            raise he

        # Step 3: Detect file extension
        ext = os.path.splitext(temp_path)[1].lower().lstrip('.')
        log_service_event("file_type_detected", "Detected file type", {
            "input": original_input,
            "local_path": temp_path,
            "extension": ext or "(none)",
            "zip_depth": _current_zip_depth
        })

        # Step 4: Handle ZIP files with depth checking
        if ext == 'zip':
            chunks = handle_zip_file(temp_path, source_url or original_input, _current_zip_depth)
            
            # Check if ZIP processing found any valid content
            if not chunks:
                log_service_event("zip_no_valid_content", "ZIP contained no valid processable content", {
                    "input": original_input,
                    "zip_depth": _current_zip_depth
                })
                # Return empty list to trigger "Files Not found" response
                return []
            
            return chunks

        # Step 5: Route to appropriate parser
        parser_func = Parser_Registry.get(ext)
        if not parser_func:
            log_service_event("unsupported_file_type", "Unsupported file type encountered", {
                "extension": ext,
                "path": temp_path,
                "zip_depth": _current_zip_depth
            })
            # Return empty list for unsupported file types
            return []

        # Step 6: Execute parser and collect chunks
        chunks = parser_func(temp_path, source_url or original_input)
        
        # Check if parser returned any content
        if not chunks:
            log_service_event("parser_no_content", "Parser returned no content", {
                "input": original_input,
                "extension": ext,
                "parser": parser_func.__name__,
                "zip_depth": _current_zip_depth
            })
            # Return empty list to trigger "Files Not found" response
            return []
        
        log_service_event("ingest_complete", "Completed ingestion successfully", {
            "input": original_input,
            "extension": ext,
            "chunks": len(chunks),
            "zip_depth": _current_zip_depth,
            "elapsed": time.time() - start
        })
        return chunks
        
    except HTTPException:
        # Re-raise HTTP exceptions (security violations, etc.)
        raise
    except Exception as e:
        # Log unexpected errors and convert to HTTP exception
        log_error("ingest_unexpected_error", {
            "input": original_input,
            "zip_depth": _current_zip_depth,
            "error": str(e),
            "error_type": type(e).__name__
        })
        raise HTTPException(
            status_code=500,
            detail=f"Unexpected error during file processing: {str(e)}"
        )
    finally:
        # Clean up temporary files (only for non-local, non-recursive calls)
        if temp_path and os.path.exists(temp_path) and not is_local and not _is_recursive:
            try:
                os.unlink(temp_path)
                log_service_event("temp_file_cleanup", "Temporary file cleaned up", {
                    "temp_path": temp_path
                })
            except OSError as e:
                log_error("temp_file_cleanup_failed", {
                    "temp_path": temp_path,
                    "error": str(e)
                })


__all__ = ["ingest_from_url", "Parser_Registry", "parse_pdf_file"]
